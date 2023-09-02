# bot.py
import os
from discord.ext import commands
import discord
import requests
import shutil
import datetime as DT
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
#TODO
#Make a dictionary of people and number of images sent
week_ago = DT.datetime.utcnow()-DT.timedelta(days=7)
# You enter the channel name here
ChannelName = 'everyday-snaps'
SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_TITLE_ONLY = 5

TOKEN = str(os.getenv('DISCORD_TOKEN'))
GUILD = os.getenv('DISCORD_GUILD')

client = discord.Client()

@client.event
async def on_ready():
    for guild in client.guilds:
        if guild.name == GUILD:
            break

    print(
        f'{client.user} is connected to the following guild:\n'
        f'{guild.name}(id: {guild.id})\n'
    )

    members = '\n - '.join([member.name for member in guild.members])
    channel = discord.utils.get(client.get_all_channels(), name=ChannelName)
    messages = await channel.history(after=week_ago, oldest_first=True).flatten()
    dict = {}
    counter = 0
    for i in messages:
        #print(i.author.name)
        #try:
        #    url = i.attachments[0].url
        #except IndexError:
        #    print("This message has no attachment.")
        #else:
            for pee in i.attachments:
                url = i.attachments[counter].url
                userName = i.author.name
                if userName not in dict.keys():
                    dict.update({userName : 0})
                else:
                    dict.update(({userName : dict.get(userName) + 1}))
                if url[0:26] == "https://cdn.discordapp.com":
                    r = requests.get(url, stream=True)
                    imageName = "Images/" + userName + str(dict.get(userName)) + '.png'
                    with open(imageName, 'wb') as outfile:
                        print("Saving image: " + imageName)
                        shutil.copyfileobj(r.raw, outfile)
                        print("Image saved!")
                counter = counter + 1
            counter = 0
    channel = discord.utils.get(client.get_all_channels(), name=ChannelName)
    print("Done!")
    print()



    """
    print("Creating PowerPoint...")
    prs = Presentation()
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Sunday Snaps!"
    slide.shapes[0].text_frame.paragraphs[0].font.name = 'Aldhabi'
    slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(88)
    time = DT.datetime.utcnow()-DT.timedelta(hours=4)
    slide.shapes[1].text_frame.paragraphs[0].text = "Created: " + str(time.month) + '/' + str(time.day) + '/' + str(time.year) + ' ' + str(time.hour) + ':' + str(time.minute) + ':' + str(time.second)
    slide.shapes[1].text_frame.paragraphs[0].font.name = 'Aldhabi'
    slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(44)

    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_ONLY]
    dir_path = os.path.dirname(os.path.realpath(__file__))

    for i in dict:
        counter = 1
        for root, dirs, files in os.walk(dir_path):
            for file in files:

                # change the extension from '.mp3' to
                # the one of your choice.
                if i in file:
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes[0].text_frame.paragraphs[0].text = i + ' (' + str(counter) + ' of ' + str(dict.get(i) + 1) + ')'
                    slide.shapes[0].text_frame.paragraphs[0].font.name = 'Aldhabi'
                    slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(44)
                    slide.shapes.add_picture(root + '/' + str(file), Inches(1.28), Inches(1.55), width=Inches(7.69), height=Inches(5.39))
                    counter = counter + 1


    prs.save("Sunday Snaps.pptx")
    print("PowerPoint Saved!")
    print()
"""
    print("Creating Excel file...")
    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'Names'
    ws['B1'] = 'Points'
    ws['D1'] = 'Team 1'
    ws['E1'] = 'Team 2'
    ws['F1'] = 'Team 3'
    ws['G1'] = 'Team 4'
    Team1, Team2, Team3, Team4 = [], [], [], []

    TeamFile = open('Team1.txt', 'r')
    Team1Names = TeamFile.readlines()
    TeamFile.close()

    TeamFile = open('Team2.txt', 'r')
    Team2Names = TeamFile.readlines()
    TeamFile.close()

    TeamFile = open('Team3.txt', 'r')
    Team3Names = TeamFile.readlines()
    TeamFile.close()

    TeamFile = open('Team4.txt', 'r')
    Team4Names = TeamFile.readlines()
    TeamFile.close()

    counter = 2
    for i in dict:
        ws['A' + str(counter)] = i
        ws['B' + str(counter)] = dict.get(i) + 1
        if i in Team1Names:
            Team1.append(dict.get(i) + 1)
        elif i in Team2Names:
            Team2.append(dict.get(i) + 1)
        elif i in Team3Names:
            Team3.append(dict.get(i) + 1)
        elif i in Team4Names:
            Team4.append(dict.get(i) + 1)

        counter = counter + 1
    ws['D2'] = sum(Team1)
    ws['E2'] = sum(Team2)
    ws['F2'] = sum(Team3)
    ws['G2'] = sum(Team4)
    wb.save("Sunday Snaps.xlsx")
    print("Excel file Saved!")
    print("This means no errors")
    exit(0)
    #await channel.send("I can now save every image sent in a channel within the past week!")



client.run("OTk5MTkxOTAxMjg4Nzk2MjYw.GsMUpi.qObMTbiNvVfCFTFyejXWNjjz0pUJxYa5b_-AL0")
#history(around=random_time)